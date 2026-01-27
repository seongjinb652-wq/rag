#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import logging
import re
import zlib
import fitz  # PyMuPDF
from pathlib import Path
from datetime import datetime
from docx import Document
from pptx import Presentation
import olefile

# 1. 경로 설정
SOURCE_DIR = Path(r"C:/Users/USER/Downloads/@@@인도네시아PDT암센터FS")
OUTPUT_DIR = Path(r"C:/Users/USER/rag/src/data/text_converted")

# 로그 설정
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("conversion_log_v2.txt", encoding='utf-8'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

class DocumentConverterV2:
    def __init__(self):
        if not OUTPUT_DIR.exists():
            OUTPUT_DIR.mkdir(parents=True)
            logger.info(f"📂 출력 디렉토리 확인/생성 완료: {OUTPUT_DIR}")

    def extract_hwp_text(self, file_path: Path) -> str:
        """HWP 파일 텍스트 추출 및 바이너리 노이즈 제거"""
        text = ""
        try:
            f = olefile.OleFileIO(str(file_path))
            dirs = f.listdir()
            bodytext = [d for d in dirs if d[0].startswith("BodyText")]
            for section in bodytext:
                data = f.openstream(section).read()
                try:
                    decompressed = zlib.decompress(data, -15)
                except:
                    decompressed = data
                text += decompressed.decode('utf-16', errors='ignore')
            f.close()
            # 노이즈 제거 필터 (한글, 영문, 숫자, 주요 문장부호 유지)
            text = re.sub(r'[^\w\s\.\,\?\!\(\)\[\]\%\:\-\d\uAC00-\uD7A3]+', ' ', text)
        except Exception as e:
            logger.error(f"❌ HWP 추출 실패 ({file_path.name}): {e}")
        return text

    def clean_filename(self, stem: str) -> str:
        """파일명에서 특수문자를 제거하여 안전한 이름 생성"""
        clean = re.sub(r'[^\w\s-]', '', stem).strip()
        return clean[:80]  # 너무 긴 파일명 방지

    def convert(self):
        target_exts = {'.pdf', '.docx', '.pptx', '.txt', '.hwp'}
        all_files = [p for p in SOURCE_DIR.rglob('*') if p.suffix.lower() in target_exts]
        
        logger.info(f"🔍 총 {len(all_files)}개 파일 변환 프로세스 시작 (V2)")

        success_count = 0
        for idx, file_path in enumerate(all_files, 1):
            try:
                ext = file_path.suffix.lower()
                content = ""

                # 확장자별 추출 로직
                if ext == '.pdf':
                    doc = fitz.open(file_path)
                    content = " ".join([page.get_text() for page in doc])
                    doc.close()
                elif ext == '.docx':
                    doc = Document(file_path)
                    content = "\n".join([p.text for p in doc.paragraphs])
                elif ext == '.hwp':
                    content = self.extract_hwp_text(file_path)
                elif ext == '.pptx':
                    prs = Presentation(file_path)
                    content = " ".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
                elif ext == '.txt':
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()

                if content.strip():
                    # 1. PDT 약어 보정: 검색 최적화를 위해 용어 병기
                    # PDT라는 단어가 단독으로 나올 때 (광역동 치료)를 추가
                    content = re.sub(r'\bPDT\b', 'PDT(광역동 치료)', content)
                    
                    # 2. 중복 단어 보정 (신축사업축사업 등)
                    content = re.sub(r'([가-힣]{2,})\1', r'\1', content)
                    
                    # 3. 파일명 정리 (해시 제거, 원본 이름 유지)
                    safe_stem = self.clean_filename(file_path.stem)
                    output_path = OUTPUT_DIR / f"{safe_stem}.txt"
                    
                    # 파일명 중복 방지 처리 (동일 이름 있을 경우 숫자 추가)
                    counter = 1
                    while output_path.exists():
                        output_path = OUTPUT_DIR / f"{safe_stem}_{counter}.txt"
                        counter += 1

                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(f"Source_Path: {file_path}\n")
                        f.write(f"Conversion_Date: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
                        f.write("-" * 50 + "\n")
                        f.write(content)
                    
                    success_count += 1
                    if idx % 10 == 0 or idx == len(all_files):
                        logger.info(f"⏳ 진행 중: [{idx}/{len(all_files)}] '{file_path.name}' 변환 완료")
                
            except Exception as e:
                logger.error(f"❌ 변환 에러 ({file_path.name}): {e}")

        logger.info(f"🏁 V2 변환 완료! 성공: {success_count}/{len(all_files)}")

if __name__ == "__main__":
    converter = DocumentConverterV2()
    converter.convert()
