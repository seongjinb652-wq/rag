#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import logging
import hashlib
import re
import zlib
import fitz  # PyMuPDF
from pathlib import Path
from datetime import datetime
from typing import List, Dict

# 라이브러리 로드
import chromadb
from openai import OpenAI
from docx import Document
from pptx import Presentation
import olefile  # pip install olefile 필수

# 1. 환경 설정
TARGET_DIR = Path(r"C:/Users/USER/Downloads/@@@인도네시아PDT암센터FS")
DB_PATH = Path(r"C:/Users/USER/rag/src/data/chroma_db")
COLLECTION_NAME = "indonesia_pdt_docs"

# OpenAI 설정
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY") 
EMBED_MODEL = "text-embedding-3-small"

CHUNK_SIZE = 800
CHUNK_OVERLAP = 100

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class RobustRAGLoaderV1:
    def __init__(self):
        logger.info("🚀 시스템 초기화 (OpenAI Embedding Mode + HWP Support)")
        self.client = OpenAI(api_key=OPENAI_API_KEY)
        
        # ChromaDB 초기화
        self.db_client = chromadb.PersistentClient(path=str(DB_PATH))
        
        # 이어서 하기(Append)를 위해 get_or_create 사용
        self.collection = self.db_client.get_or_create_collection(
            name=COLLECTION_NAME,
            metadata={"hnsw:space": "cosine"}
        )
        
    def extract_hwp_text(self, file_path: Path) -> str:
        import olefile
        import zlib
        import re

        text = ""
        try:
            f = olefile.OleFileIO(str(file_path))
            dirs = f.listdir()
            if ["FileHeader"] not in dirs: return ""

            bodytext = [d for d in dirs if d[0].startswith("BodyText")]
            for section in bodytext:
                data = f.openstream(section).read()
                try:
                    decompressed = zlib.decompress(data, -15)
                except:
                    decompressed = data
                
                # 깨끗한 텍스트만 추출하기 위한 처리
                raw_content = decompressed.decode('utf-16', errors='ignore')
                
                # [핵심] 한글, 영문, 숫자, 일반 문장부호만 남기고 제거
                clean_content = re.sub(r'[^\w\s\.\,\?\!\(\)\[\]\%\:\-\d\uAC00-\uD7A3]+', '', raw_content)
                text += clean_content
            f.close()
        except Exception as e:
            logger.error(f"❌ {file_path.name} HWP 추출 실패: {e}")
        return text
    

    def extract_text(self, file_path: Path) -> str:
        """파일별 텍스트 추출 및 보정"""
        ext = file_path.suffix.lower()
        text = ""
        try:
            if ext == '.pdf':
                doc = fitz.open(file_path)
                text = " ".join([page.get_text() for page in doc])
                doc.close()
            elif ext == '.docx':
                doc = Document(file_path)
                text = "\n".join([p.text for p in doc.paragraphs])
            elif ext == '.hwp':
                text = self.extract_hwp_text(file_path)
            elif ext == '.pptx':
                prs = Presentation(file_path)
                text = " ".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
            elif ext == '.txt':
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read()

            if text:
                # 1. 한글 인코딩 오류 교정
                corrections = {
                    "싞": "신", "짂": "진", "읶": "인", "핚": "한", 
                    "웎": "원", "상홖": "상환", "읷": "일", "젂": "전",
                    "곾": "관", "첚": "천", "중갂": "중간", "얶": "언",
                    "읻": "인", "싵": "실", "엓": "엑", "짗": "지",
                    "읃": "은", "읔": "억", "싲": "시", "짘": "직"
                }
                for wrong, right in corrections.items():
                    text = text.replace(wrong, right)
                
                # 2. 중복 단어 제거
                text = re.sub(r'([가-힣]{2,})\1', r'\1', text)

        except Exception as e:
            logger.error(f"❌ {file_path.name} 처리 에러: {e}")
        
        return text

    def run(self):
        # 1. 모든 파일 스캔 (HWP 추가)
        all_files = []
        target_exts = {'.pdf', '.docx', '.pptx', '.txt', '.hwp'}
        for root, _, filenames in os.walk(TARGET_DIR):
            for f in filenames:
                p = Path(root) / f
                if p.suffix.lower() in target_exts:
                    all_files.append(p)

        logger.info(f"🔍 총 {len(all_files)}개 파일 발견")

        # 2. 이미 DB에 저장된 파일 확인
        existing_data = self.collection.get(include=['metadatas'])
        processed_files = {m['source'] for m in existing_data['metadatas']} if existing_data['metadatas'] else set()

        total_chunks = 0
        for idx, file_path in enumerate(all_files, 1):
            if file_path.name in processed_files:
                continue

            raw_text = self.extract_text(file_path)
            if not raw_text.strip():
                continue

            chunks = [raw_text[i:i+CHUNK_SIZE] for i in range(0, len(raw_text), CHUNK_SIZE - CHUNK_OVERLAP)]
            
            try:
                response = self.client.embeddings.create(input=chunks, model=EMBED_MODEL)
                embeddings = [data.embedding for data in response.data]
                
                ids = [hashlib.md5(f"{file_path.name}_{i}".encode()).hexdigest() for i in range(len(chunks))]
                self.collection.add(
                    ids=ids,
                    embeddings=embeddings,
                    documents=chunks,
                    metadatas=[{"source": file_path.name} for _ in chunks]
                )
                total_chunks += len(chunks)
                logger.info(f"✅ [{idx}/{len(all_files)}] {file_path.name} ({len(chunks)} Chunks)")
            
            except Exception as e:
                logger.error(f"❌ [{idx}/{len(all_files)}] {file_path.name} 임베딩 에러: {e}")

        logger.info(f"🏁 완료! 현재 컬렉션 내 총 청크 수: {self.collection.count()}")

if __name__ == "__main__":
    loader = RobustRAGLoaderV1()
    loader.run()
