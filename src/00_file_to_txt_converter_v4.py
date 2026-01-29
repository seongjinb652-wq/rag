import os
import logging
import re
import zlib
import hashlib
import pdfplumber  # pip install pdfplumber
import olefile
from pathlib import Path
from datetime import datetime
from docx import Document
from pptx import Presentation

# 1. 경로 설정
SOURCE_DIR = Path(r"C:/Users/USER/Downloads/@@@인도네시아PDT암센터FS")
OUTPUT_DIR = Path(r"C:/Users/USER/rag/src/data/text_converted")

# 로그 설정
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("conversion_log_v4.txt", encoding='utf-8'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

class DocumentConverterV4:
    """
    [출처보완] 대용량 증분 처리 및 메타데이터 정밀 추출 버전
    """
    def __init__(self):
        if not OUTPUT_DIR.exists():
            OUTPUT_DIR.mkdir(parents=True)
            logger.info(f"📂 v4 출력 디렉토리 생성: {OUTPUT_DIR}")

    def format_as_markdown(self, table):
        """표 데이터를 마크다운으로 변환"""
        if not table or not any(any(cell for cell in row) for row in table):
            return ""
        
        markdown = "\n"
        for i, row in enumerate(table):
            clean_row = [str(cell).replace('\n', ' ').strip() if cell is not None else "" for cell in row]
            markdown += "| " + " | ".join(clean_row) + " |\n"
            if i == 0:
                markdown += "| " + " | ".join(["---"] * len(row)) + " |\n"
        return markdown + "\n"

    def extract_pdf_smart(self, file_path):
        """PDF 텍스트 및 표 추출"""
        full_content = []
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text() or ""
                    full_content.append(text)
                    tables = page.extract_tables()
                    for table in tables:
                        md_table = self.format_as_markdown(table)
                        if md_table.strip():
                            full_content.append(md_table)
            return "\n".join(full_content)
        except Exception as e:
            logger.error(f"❌ PDF 추출 실패: {e}")
            return ""

    def extract_hwp_text(self, file_path):
        """HWP 텍스트 추출 (v4 깨짐 방지 로직)"""
        try:
            f = olefile.OleFileIO(str(file_path))
            dirs = f.listdir()
            bodytext = [d for d in dirs if d[0].startswith("BodyText/Section")]
            if not bodytext: bodytext = [d for d in dirs if d[0].startswith("BodyText")]
                
            full_text = ""
            for section in bodytext:
                data = f.openstream(section).read()
                try:
                    decompressed = zlib.decompress(data, -15)
                except:
                    decompressed = data
                section_text = decompressed.decode('utf-16', errors='ignore')
                full_text += "".join([c for c in section_text if ord(c) >= 32 or c in "\n\r\t"])
            
            f.close()
            return re.sub(r'[^\w\s\.\,\?\!\(\)\[\]\%\:\-\d\uAC00-\uD7A3]+', ' ', full_text)
        except Exception as e:
            logger.error(f"❌ HWP 추출 실패: {e}")
            return ""

    def convert(self):
        target_exts = {'.pdf', '.docx', '.pptx', '.txt', '.hwp'}
        # 전체 경로 탐색
        all_files = [p for p in SOURCE_DIR.rglob('*') if p.suffix.lower() in target_exts]
        
        logger.info(f"🚀 [v4 출처보완] 총 {len(all_files)}개 파일 변환 시작...")

        success_count = 0
        for idx, file_path in enumerate(all_files, 1):
            try:
                if file_path.name.startswith("~$"): continue

                ext = file_path.suffix.lower()
                content = ""

                # 파일 타입별 추출
                if ext == '.pdf': content = self.extract_pdf_smart(file_path)
                elif ext == '.docx':
                    doc = Document(file_path)
                    content = "\n".join([p.text for p in doc.paragraphs])
                elif ext == '.hwp': content = self.extract_hwp_text(file_path)
                elif ext == '.pptx':
                    prs = Presentation(file_path)
                    content = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
                elif ext == '.txt':
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()

                if content.strip():
                    # 1. 용어 보정
                    content = re.sub(r'\bPDT\b', 'PDT(광역동 치료)', content)
                    
                    # 2. 경로 표준화 (Windows 역슬래시를 슬래시로 통일하여 저장)
                    standard_path = str(file_path.absolute()).replace('\\', '/')
                    
                    # 3. 고유 파일명 생성 (파일명 + 경로 해시 조합)
                    path_hash = hashlib.md5(standard_path.encode()).hexdigest()[:8]
                    safe_stem = re.sub(r'[^\w\s-]', '', file_path.stem).strip()[:40]
                    output_filename = f"{safe_stem}_{path_hash}.txt"
                    output_path = OUTPUT_DIR / output_filename
                    
                    # 4. 파일 저장 (나중에 읽기 쉬운 표준 포맷)
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(f"Source: {standard_path}\n")
                        f.write("-" * 60 + "\n")
                        f.write(content)
                    
                    success_count += 1
                    if idx % 10 == 0 or idx == len(all_files):
                        logger.info(f"⏳ 진행 중: [{idx}/{len(all_files)}] {file_path.name} 완료")
                
            except Exception as e:
                logger.error(f"❌ 변환 에러 ({file_path.name}): {e}")

        logger.info(f"🏁 v4 변환 완료! (성공: {success_count}/{len(all_files)})")

if __name__ == "__main__":
    converter = DocumentConverterV4()
    converter.convert()
