#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import logging
import re
import zlib
import fitz  # PyMuPDF
from pathlib import Path
from docx import Document
from pptx import Presentation
import olefile

# 1. 경로 설정
SOURCE_DIR = Path(r"C:/Users/USER/Downloads/@@@인도네시아PDT암센터FS")
OUTPUT_DIR = Path(r"C:/Users/USER/rag/src/data/text_converted")

# 로그 설정
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("conversion_log.txt", encoding='utf-8'),
        logging.Stream_StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

class DocumentConverter:
    def __init__(self):
        if not OUTPUT_DIR.exists():
            OUTPUT_DIR.mkdir(parents=True)
            logger.info(f"📂 출력 디렉토리 생성 완료: {OUTPUT_DIR}")

    def extract_hwp_text(self, file_path: Path) -> str:
        text = ""
        try:
            f = olefile.OleFileIO(str(file_path))
            dirs = f.listdir()
            bodytext = [d for d in dirs if d[0].startswith("BodyText")]
            for section in bodytext:
                data = f.openstream(section).read()
                try:
                    decompressed = zlib.decompress(data, -15)
                except:
                    decompressed = data
                text += decompressed.decode('utf-16', errors='ignore')
            f.close()
            # 바이너리 노이즈 제거 필터
            text = re.sub(r'[^\w\s\.\,\?\!\(\)\[\]\%\:\-\d\uAC00-\uD7A3]+', ' ', text)
        except Exception as e:
            logger.error(f"❌ HWP 추출 실패 ({file_path.name}): {e}")
        return text

    def convert(self):
        target_exts = {'.pdf', '.docx', '.pptx', '.txt', '.hwp'}
        all_files = [p for p in SOURCE_DIR.rglob('*') if p.suffix.lower() in target_exts]
        
        logger.info(f"🔍 총 {len(all_files)}개 파일 변환 시작...")

        success_count = 0
        for idx, file_path in enumerate(all_files, 1):
            try:
                ext = file_path.suffix.lower()
                content = ""

                if ext == '.pdf':
                    doc = fitz.open(file_path)
                    content = " ".join([page.get_text() for page in doc])
                    doc.close()
                elif ext == '.docx':
                    doc = Document(file_path)
                    content = "\n".join([p.text for p in doc.paragraphs])
                elif ext == '.hwp':
                    content = self.extract_hwp_text(file_path)
                elif ext == '.pptx':
                    prs = Presentation(file_path)
                    content = " ".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
                elif ext == '.txt':
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()

                if content.strip():
                    # 중복 단어 보정 로직 적용
                    content = re.sub(r'([가-힣]{2,})\1', r'\1', content)
                    
                    # 출력 파일 경로 설정 (원본 구조 유지 대신 평면 저장 혹은 구조 재현 가능)
                    # 여기서는 찾기 쉽게 파일명.txt로 저장합니다.
                    safe_name = f"{file_path.stem}_{hashlib.md5(str(file_path).encode()).hexdigest()[:6]}.txt"
                    output_path = OUTPUT_DIR / safe_name
                    
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(f"Source: {file_path}\n")
                        f.write("-" * 50 + "\n")
                        f.write(content)
                    
                    success_count += 1
                    if idx % 10 == 0:
                        logger.info(f"⏳ 진행 중: [{idx}/{len(all_files)}] 변환 완료")
                
            except Exception as e:
                logger.error(f"❌ 변환 에러 ({file_path.name}): {e}")

        logger.info(f"🏁 변환 완료! 성공: {success_count}/{len(all_files)}")
        logger.info(f"📂 저장 위치: {OUTPUT_DIR}")

import hashlib
if __name__ == "__main__":
    converter = DocumentConverter()
    converter.convert()
