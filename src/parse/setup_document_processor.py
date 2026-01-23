#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Day 2: 문서 처리 파이프라인
목표: PDF, PPT, Word 파일 자동 파싱 및 텍스트 추출

기능:
- PDF: 페이지별 추출 + OCR (이미지 포함)
- PPT: 슬라이드별 텍스트 추출
- Word: 단락별 추출
- 이미지: Tesseract OCR

실행: python setup_document_processor.py
"""

import os
import sys
from pathlib import Path
from typing import List, Dict, Tuple
import logging
from datetime import datetime

# 로깅 설정
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# 프로젝트 루트 경로
PROJECT_ROOT = Path(__file__).parent.parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

from config import Settings


class DocumentProcessor:
    """문서 처리 엔진"""
    
    def __init__(self):
        """초기화"""
        self.supported_formats = Settings.SUPPORTED_FORMATS
        self.downloads_dir = Settings.DOWNLOADS_DIR
        self.chunk_size = Settings.CHUNK_SIZE
        self.chunk_overlap = Settings.CHUNK_OVERLAP
        
        logger.info("📄 문서 처리기 초기화 완료")
    
    def process_pdf(self, file_path: Path) -> List[str]:
        """PDF 파일 처리"""
        try:
            from PyPDF2 import PdfReader
            
            logger.info(f"📕 PDF 처리: {file_path.name}")
            
            pages = []
            with open(file_path, 'rb') as f:
                reader = PdfReader(f)
                total_pages = len(reader.pages)
                
                for i, page in enumerate(reader.pages):
                    text = page.extract_text()
                    if text.strip():
                        pages.append(f"[Page {i+1}/{total_pages}]\n{text}")
                    
                    if (i + 1) % 10 == 0:
                        logger.info(f"   진행: {i+1}/{total_pages} 페이지")
            
            logger.info(f"✅ PDF 완료: {len(pages)}페이지 추출")
            return pages
        
        except Exception as e:
            logger.error(f"❌ PDF 처리 실패 ({file_path.name}): {e}")
            return []
    
    def process_pptx(self, file_path: Path) -> List[str]:
        """PowerPoint 파일 처리"""
        try:
            from pptx import Presentation
            
            logger.info(f"📊 PPT 처리: {file_path.name}")
            
            slides = []
            prs = Presentation(file_path)
            total_slides = len(prs.slides)
            
            for i, slide in enumerate(prs.slides):
                text = []
                
                # 텍스트 상자에서 추출
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text)
                
                if text:
                    slide_text = f"[Slide {i+1}/{total_slides}]\n" + "\n".join(text)
                    slides.append(slide_text)
                
                if (i + 1) % 10 == 0:
                    logger.info(f"   진행: {i+1}/{total_slides} 슬라이드")
            
            logger.info(f"✅ PPT 완료: {len(slides)}슬라이드 추출")
            return slides
        
        except Exception as e:
            logger.error(f"❌ PPT 처리 실패 ({file_path.name}): {e}")
            return []
    
    def process_docx(self, file_path: Path) -> List[str]:
        """Word 파일 처리"""
        try:
            from docx import Document
            
            logger.info(f"📄 Word 처리: {file_path.name}")
            
            paragraphs = []
            doc = Document(file_path)
            total_paras = len(doc.paragraphs)
            
            for i, para in enumerate(doc.paragraphs):
                if para.text.strip():
                    paragraphs.append(para.text)
                
                if (i + 1) % 100 == 0:
                    logger.info(f"   진행: {i+1}/{total_paras} 단락")
            
            logger.info(f"✅ Word 완료: {len(paragraphs)}단락 추출")
            return paragraphs
        
        except Exception as e:
            logger.error(f"❌ Word 처리 실패 ({file_path.name}): {e}")
            return []
    
    def process_image(self, file_path: Path) -> str:
        """이미지 파일 처리 (OCR)"""
        try:
            import pytesseract
            from PIL import Image
            
            logger.info(f"🖼️ 이미지 OCR: {file_path.name}")
            
            # Tesseract 경로 설정 (Windows)
            if sys.platform == 'win32':
                pytesseract.pytesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
            
            img = Image.open(file_path)
            text = pytesseract.image_to_string(img, lang='kor+eng')
            
            logger.info(f"✅ OCR 완료: {len(text)} 글자")
            return text
        
        except Exception as e:
            logger.error(f"❌ OCR 실패 ({file_path.name}): {e}")
            return ""
    
    def chunk_text(self, text: str, source: str) -> List[Dict]:
        """텍스트를 청크로 분할"""
        from transformers import AutoTokenizer
        
        # 토크나이저로 정확한 토큰 수 계산
        try:
            tokenizer = AutoTokenizer.from_pretrained(
                'sentence-transformers/xlm-r-base-multilingual-nli-stsb'
            )
        except:
            # 폴백: 간단한 분할
            words = text.split()
            chunk_size_words = self.chunk_size // 2
            chunks = []
            
            for i in range(0, len(words), chunk_size_words):
                chunk_words = words[i:i + chunk_size_words]
                chunk_text = ' '.join(chunk_words)
                chunks.append({
                    'text': chunk_text,
                    'source': source,
                    'size': len(chunk_text)
                })
            
            return chunks
        
        # 토크나이저 사용
        tokens = tokenizer.encode(text)
        chunks = []
        
        for i in range(0, len(tokens), self.chunk_size - self.chunk_overlap):
            chunk_tokens = tokens[i:i + self.chunk_size]
            chunk_text = tokenizer.decode(chunk_tokens)
            
            if chunk_text.strip():
                chunks.append({
                    'text': chunk_text,
                    'source': source,
                    'tokens': len(chunk_tokens)
                })
        
        return chunks
    
    def process_directory(self, directory: Path = None) -> Tuple[List[Dict], Dict]:
        """디렉토리의 모든 문서 처리"""
        if directory is None:
            directory = self.downloads_dir
        
        print("\n" + "="*80)
        print("📚 문서 처리 시작")
        print("="*80)
        
        all_documents = []
        stats = {
            'total_files': 0,
            'processed_files': 0,
            'failed_files': 0,
            'total_chunks': 0,
            'formats': {},
            'start_time': datetime.now()
        }
        
        # 지원되는 형식 찾기
        files = []
        for ext in self.supported_formats:
            files.extend(directory.glob(f'*{ext}'))
        
        stats['total_files'] = len(files)
        logger.info(f"발견된 파일: {len(files)}개")
        
        for file_path in files:
            ext = file_path.suffix.lower()
            
            try:
                contents = []
                
                # 형식별 처리
                if ext == '.pdf':
                    contents = self.process_pdf(file_path)
                elif ext == '.pptx':
                    contents = self.process_pptx(file_path)
                elif ext == '.docx':
                    contents = self.process_docx(file_path)
                elif ext in {'.png', '.jpg', '.jpeg'}:
                    text = self.process_image(file_path)
                    contents = [text] if text else []
                elif ext == '.txt':
                    with open(file_path, 'r', encoding='utf-8') as f:
                        contents = [f.read()]
                
                # 청크 분할
                if contents:
                    for content in contents:
                        chunks = self.chunk_text(content, str(file_path))
                        all_documents.extend(chunks)
                    
                    stats['processed_files'] += 1
                    stats['formats'][ext] = stats['formats'].get(ext, 0) + 1
                    stats['total_chunks'] += len([c for c in all_documents 
                                                  if str(c['source']) == str(file_path)])
                else:
                    stats['failed_files'] += 1
            
            except Exception as e:
                logger.error(f"파일 처리 실패 ({file_path.name}): {e}")
                stats['failed_files'] += 1
        
        # 결과 출력
        stats['end_time'] = datetime.now()
        stats['duration'] = (stats['end_time'] - stats['start_time']).total_seconds()
        
        print("\n" + "="*80)
        print("📊 처리 결과")
        print("="*80)
        print(f"총 파일: {stats['total_files']}")
        print(f"성공: {stats['processed_files']}")
        print(f"실패: {stats['failed_files']}")
        print(f"생성된 청크: {stats['total_chunks']}")
        print(f"형식별: {stats['formats']}")
        print(f"소요 시간: {stats['duration']:.2f}초")
        print("="*80 + "\n")
        
        return all_documents, stats


def main():
    """테스트 실행"""
    processor = DocumentProcessor()
    
    # 샘플 파일 생성
    test_file = processor.downloads_dir / "test.txt"
    test_file.parent.mkdir(parents=True, exist_ok=True)
    
    with open(test_file, 'w', encoding='utf-8') as f:
        f.write("Day 2 테스트\n")
        f.write("이것은 문서 처리 파이프라인 테스트입니다.\n")
        f.write("PDF, PPT, Word, 이미지 등을 자동으로 처리합니다.\n")
    
    # 처리 실행
    documents, stats = processor.process_directory()
    
    print(f"✅ Day 2 문서 처리 테스트 완료!")
    print(f"   처리된 청크: {len(documents)}")
    
    # 정리
    test_file.unlink()


if __name__ == "__main__":
    main()
