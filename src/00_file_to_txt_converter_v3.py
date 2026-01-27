# 마크다운 적용 버전
import os
import logging
import re
import zlib
import hashlib
import pdfplumber  # pip install pdfplumber
from pathlib import Path
from datetime import datetime
from docx import Document
from pptx import Presentation
import olefile

# 1. 경로 설정
SOURCE_DIR = Path(r"C:/Users/USER/Downloads/@@@인도네시아PDT암센터FS")
OUTPUT_DIR = Path(r"C:/Users/USER/rag/src/data/text_converted")

# 로그 설정
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("conversion_log_v3.txt", encoding='utf-8'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

class DocumentConverterV3:
    def __init__(self):
        if not OUTPUT_DIR.exists():
            OUTPUT_DIR.mkdir(parents=True)
            logger.info(f"📂 출력 디렉토리 확인: {OUTPUT_DIR}")

    def format_as_markdown(self, table):
        """리스트 형태의 표 데이터를 마크다운 문자열로 변환"""
        if not table or not any(any(cell for cell in row) for row in table):
            return ""
        
        markdown = "\n"
        for i, row in enumerate(table):
            # 셀 내 줄바꿈 제거, None 처리, 불필요한 공백 제거
            clean_row = [str(cell).replace('\n', ' ').strip() if cell is not None else "" for cell in row]
            markdown += "| " + " | ".join(clean_row) + " |\n"
            if i == 0: # 헤더와 본문 구분선
                markdown += "| " + " | ".join(["---"] * len(row)) + " |\n"
        return markdown + "\n"

    def extract_pdf_smart(self, file_path):
        """PDF에서 텍스트와 표를 마크다운 구조로 추출"""
        full_content = []
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    # 1. 텍스트 추출
                    text = page.extract_text() or ""
                    full_content.append(text)
                    
                    # 2. 표 추출 및 마크다운 변환
                    tables = page.extract_tables()
                    for table in tables:
                        md_table = self.format_as_markdown(table)
                        if md_table.strip():
                            full_content.append(md_table)
            return "\n".join(full_content)
        except Exception as e:
            logger.error(f"❌ PDF 스마트 추출 실패 ({file_path.name}): {e}")
            return ""
################################
    import olefile
    import zlib
    import struct

    def extract_hwp_text_v4(file_path):
        """HWP 파일의 깨짐을 최소화한 텍스트 추출 함수"""
        try:
            f = olefile.OleFileIO(file_path)
            dirs = f.listdir()
        
            # HWP 버전 및 암호화 여부 체크 (FileHeader)
            if ["FileHeader"] not in dirs:
                return ""
            
            # 본문 데이터가 담긴 Section들 찾기
            bodytext = [d for d in dirs if d[0].startswith("BodyText/Section")]
            full_text = ""
        
            for section in bodytext:
                data = f.openstream(section).read()
                # 압축 여부 확인 (HWP는 보통 압축되어 있음)
                try:
                    # zlib 압축 해제 (-15는 헤더 없는 raw deflate 대응)
                    decompressed = zlib.decompress(data, -15)
                except:
                    decompressed = data
            
                # 한글(utf-16) 디코딩 및 제어문자 정제
                section_text = decompressed.decode('utf-16', errors='ignore')
            
                # HWP 특유의 제어문자(글자 크기, 글꼴 변경 등) 제거
                clean_text = ""
                for char in section_text:
                    if ord(char) >= 32 or char in "\n\r\t":
                        clean_text += char
                full_text += clean_text
            
            f.close()
            return full_text
        except Exception as e:
            return f"Error: {str(e)}"

    
#    def extract_hwp_text(self, file_path: Path) -> str:
#        """HWP 파일 텍스트 추출"""
#        text = ""
#        try:
#            f = olefile.OleFileIO(str(file_path))
#            dirs = f.listdir()
#            bodytext = [d for d in dirs if d[0].startswith("BodyText")]
#            for section in bodytext:
#                data = f.openstream(section).read()
#                try: decompressed = zlib.decompress(data, -15)
#                except: decompressed = data
#                text += decompressed.decode('utf-16', errors='ignore')
#            f.close()
#            text = re.sub(r'[^\w\s\.\,\?\!\(\)\[\]\%\:\-\d\uAC00-\uD7A3]+', ' ', text)
#        except Exception as e:
#            logger.error(f"❌ HWP 추출 실패 ({file_path.name}): {e}")
#        return text

    def convert(self):
        target_exts = {'.pdf', '.docx', '.pptx', '.txt', '.hwp'}
        all_files = [p for p in SOURCE_DIR.rglob('*') if p.suffix.lower() in target_exts]
        
        logger.info(f"🚀 총 {len(all_files)}개 파일 스마트 변환(V3) 시작...")

        success_count = 0
        for idx, file_path in enumerate(all_files, 1):
            try:
                ext = file_path.suffix.lower()
                content = ""

                if ext == '.pdf':
                    content = self.extract_pdf_smart(file_path)
                elif ext == '.docx':
                    doc = Document(file_path)
                    content = "\n".join([p.text for p in doc.paragraphs])
                elif ext == '.hwp':
                    content = self.extract_hwp_text(file_path)
                elif ext == '.pptx':
                    prs = Presentation(file_path)
                    content = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
                elif ext == '.txt':
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()

                if content.strip():
                    # 용어 보정 (광역동 치료 추가)
                    content = re.sub(r'\bPDT\b', 'PDT(광역동 치료)', content)
                    
                    # 파일명 정리 및 저장
                    safe_stem = re.sub(r'[^\w\s-]', '', file_path.stem).strip()[:50]
                    output_path = OUTPUT_DIR / f"{safe_stem}.txt"
                    
                    # 중복 이름 처리
                    counter = 1
                    orig_path = output_path
                    while output_path.exists():
                        output_path = OUTPUT_DIR / f"{safe_stem}_{counter}.txt"
                        counter += 1

                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(f"Source: {file_path}\n---\n{content}")
                    
                    success_count += 1
                    if idx % 10 == 0 or idx == len(all_files):
                        logger.info(f"⏳ 진행 중: [{idx}/{len(all_files)}] {file_path.name} 완료")
                
            except Exception as e:
                logger.error(f"❌ 변환 에러 ({file_path.name}): {e}")

        logger.info(f"🏁 V3 스마트 변환 완료! 성공: {success_count}/{len(all_files)}")

if __name__ == "__main__":
    converter = DocumentConverterV3()
    converter.convert()
